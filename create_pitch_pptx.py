#!/usr/bin/env python3
"""
Generate Lilita Keper Business Pitch PowerPoint Presentation
Install: pip install python-pptx
Run: python create_pitch_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
GOLD = RGBColor(212, 175, 55)
DARK_EARTH = RGBColor(44, 36, 22)
CREAM = RGBColor(250, 248, 243)
CHARCOAL = RGBColor(58, 58, 58)
FOREST = RGBColor(47, 82, 51)
TEXT_SECONDARY = RGBColor(102, 102, 102)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_EARTH

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = CREAM

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
    footer_frame = footer_box.text_frame
    p = footer_frame.paragraphs[0]
    p.text = "40,600+ travel professionals • 3x revenue potential • 90-day launch"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SECONDARY

def add_content_slide(prs, title, content_list, bg_color=None):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if bg_color is None:
        bg_color = CREAM
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_EARTH

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = CHARCOAL if bg_color == CREAM else CHARCOAL
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.level = 0

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_EARTH

    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(4.2), Inches(0.4))
    left_title_frame = left_title_box.text_frame
    p = left_title_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_EARTH

    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(4.2), Inches(5))
    left_frame = left_content_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(4)
        p.space_after = Pt(4)

    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.4))
    right_title_frame = right_title_box.text_frame
    p = right_title_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = FOREST

    right_content_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
    right_frame = right_content_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(14)
        p.font.bold = True if "$" in item else False
        p.font.color.rgb = FOREST if "$" in item else CHARCOAL
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# SLIDE 1: Title
add_title_slide(prs,
    "Lilita Keper\nRevenue Acceleration",
    "Partnership Program for Luxury Bookings, Sales Management & International Workshops")

# SLIDE 2: The Opportunity
add_content_slide(prs, "The Opportunity", [
    "🌍 Global Reach: 40,600 vetted travel professionals ready to send clients",
    "⏰ Time Savings: Automated booking, calendar & responses—manage from anywhere",
    "💰 Revenue Multiplier: Direct bookings + travel agent commissions + premium packages",
    "📈 Scalable: Same system for 1 suite or 50—grow without adding staff",
    "🎯 Focus: We handle bookings, marketing, agents—you run the lodge"
], RGBColor(245, 241, 232))

# SLIDE 3: Current vs Potential
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = CREAM

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Current State vs. Potential (Year 1)"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = DARK_EARTH

# Left side - Current
current_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(4.2), Inches(5.5))
current_frame = current_box.text_frame
current_frame.word_wrap = True
items = [
    "TODAY",
    "",
    "5 Suites",
    "20-30% Occupancy",
    "Manual Booking",
    "Limited Reach",
    "Basic Tracking"
]
for i, item in enumerate(items):
    if i > 0:
        current_frame.add_paragraph()
    p = current_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(16)
    p.font.bold = i == 0
    p.font.color.rgb = TEXT_SECONDARY if i == 0 else CHARCOAL

# Right side - Potential
potential_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(5.5))
potential_frame = potential_box.text_frame
potential_frame.word_wrap = True
items = [
    "YEAR 1",
    "",
    "50-65% Occupancy",
    "3x Revenue",
    "Fully Automated",
    "40.6K Agent Network",
    "Real-Time Dashboard"
]
for i, item in enumerate(items):
    if i > 0:
        potential_frame.add_paragraph()
    p = potential_frame.paragraphs[i]
    p.text = item
    p.font.size = Pt(16)
    p.font.bold = i == 0 or "$" in item
    p.font.color.rgb = FOREST if i == 0 else (GOLD if "$" in item else CHARCOAL)

# SLIDE 4: How It Works - 4 Revenue Streams
add_content_slide(prs, "4 Revenue Streams", [
    "1️⃣ Direct Bookings: Guests book on your website. Est. 10-15/month at $3-5K = $180-270K",
    "",
    "2️⃣ Luxury Curator Partnerships: Luxury Curators, selective travel agents book premium packages. Est. 8-12/month = $150-200K",
    "",
    "3️⃣ Photography Tours & Content: Professional photographers, travel media, influencers + licensing. Est. 3-5/month = $80-120K",
    "",
    "4️⃣ High-Commission Travel Network: TUI partners, luxury agencies (20-25% commission). Est. 10-20/month = $100-150K"
], RGBColor(245, 241, 232))

# SLIDE 5: Platform Features
add_content_slide(prs, "Your Integrated Platform", [
    "📅 Live Booking Calendar: Real-time availability on your website. Instant confirmations.",
    "",
    "📊 Account Dashboard: Bookings, revenue, agent performance, occupancy %, analytics",
    "",
    "🤖 AI Agents: Auto-responses to inquiries, commission calculations, smart follow-ups",
    "",
    "📧 Contact Management: 40.6K contacts organized by type. Segment & campaign.",
    "",
    "💰 Finance Integration: Commission tracking, payment processing, tax-ready reports"
])

# SLIDE 6: 90-Day Launch Timeline
add_content_slide(prs, "90-Day Launch Roadmap", [
    "Week 1-2: Setup & Onboarding",
    "  • Calendly embedded on your site",
    "  • Platform trained on your pricing & policies",
    "",
    "Week 2-4: Agent Outreach Launch",
    "  • 200 travel agents contacted",
    "  • First bookings & responses",
    "",
    "Week 4-8: Scale & Optimize",
    "  • Expand to 1,000+ agents",
    "  • Launch corporate/wedding packages",
    "",
    "Week 8-12: Full Integration",
    "  • First workshop booked",
    "  • Handoff to your team"
])

# SLIDE 7: Revenue Projection
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = CREAM

title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "Year 1 Revenue Projection"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = DARK_EARTH

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(0.4))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.paragraphs[0]
p.text = "Conservative estimates based on 50-65% occupancy + partnerships"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_SECONDARY

# Content
content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(5))
text_frame = content_box.text_frame
text_frame.word_wrap = True

streams = [
    ("Direct Bookings", "$180K - 270K", "10-15 bookings/month at $3-5K"),
    ("Luxury Curator Partnerships", "$150K - 200K", "8-12 premium bookings/month @ $4-6K"),
    ("Photography Tours & Content", "$80K - 120K", "3-5 tours/month + content licensing"),
    ("High-Commission Travel Network", "$100K - 150K", "10-20 bookings/month @ 20-25%"),
    ("TOTAL POTENTIAL", "$510K - $740K", "85% net to lodge (lower take-rate fees)")
]

for i, (label, amount, desc) in enumerate(streams):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = f"{label}: {amount}"
    p.font.size = Pt(16)
    p.font.bold = i == len(streams) - 1
    p.font.color.rgb = GOLD if i == len(streams) - 1 else DARK_EARTH
    p.space_after = Pt(2)

    text_frame.add_paragraph()
    p = text_frame.paragraphs[len(text_frame.paragraphs)-1]
    p.text = f"  {desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_SECONDARY

# SLIDE 8: Investment & Contract
add_content_slide(prs, "Investment & Contract", [
    "💵 Setup Fee: $2,500 one-time (implementation, training, integration)",
    "",
    "📊 Monthly Platform: $500/month (calendar, bookings, contacts, reporting)",
    "    • Cancel anytime after 6-month minimum",
    "",
    "💰 Performance Commissions:",
    "    • Travel agent bookings: 15%",
    "    • Corporate/wedding packages: 10% markup",
    "    • Workshops: 20% facilitation fee",
    "",
    "✅ BREAK-EVEN: 4-5 additional bookings/month covers all costs",
    "   Most lodges see ROI within 30 days"
])

# SLIDE 9: What's Included
add_content_slide(prs, "What's Included", [
    "✓ Live calendar + booking system on your website",
    "✓ Automated guest confirmations & reminders",
    "✓ 40.6K travel agent database + outreach campaigns",
    "✓ Bounce recovery (auto-find new agent contacts)",
    "✓ Commission tracking & monthly reporting",
    "✓ AI booking assistant for 24/7 inquiries",
    "✓ Account management dashboard with analytics",
    "✓ Monthly strategy review + optimization",
    "✓ Email support + quarterly business reviews"
])

# SLIDE 10: Getting Started
add_content_slide(prs, "Getting Started: Next 7 Days", [
    "Today: Review this proposal with your team",
    "",
    "Day 2-3: Discovery call (30 min) review your booking process & goals",
    "",
    "Day 4-5: Custom demo of the platform pre-configured with YOUR lodge",
    "",
    "Day 6-7: Sign contract & launch - Calendly live, agents contacted"
], RGBColor(245, 241, 232))

# SLIDE 11: FAQ
add_content_slide(prs, "Common Questions", [
    "Q: What if we don't meet targets?",
    "  A: No penalties. Only pay on actual bookings. $500/month covers platform.",
    "",
    "Q: Can we cancel?",
    "  A: Yes, after 6 months. Most lodges see results in 30-60 days.",
    "",
    "Q: Will this work with our existing system?",
    "  A: Yes. Integrates with Airbnb, Booking.com, your website.",
    "",
    "Q: Who owns our contact data?",
    "  A: You do. All data stays with you, transferable if you leave."
])

# SLIDE 12: The Offer - Close
add_title_slide(prs,
    "The Offer",
    "Fully integrated booking + sales + account management system.\n\nNo upfront risk. Pay for results.\n\n30-day trial. See it work before committing.\n\nYour data is yours. Always.")

# Save presentation
output_file = "Lilita_Keper_Business_Pitch.pptx"
prs.save(output_file)
print(f"✅ PowerPoint created: {output_file}")
print("\nReady to present to lodge owners!")
